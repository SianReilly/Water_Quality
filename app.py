# pip install streamlit plotly pandas python-pptx pyarrow requests kaleido
"""
London Water Quality Explorer
==============================
Visualises Environment Agency water-quality sampling points across Greater
London and how a chosen determinand (e.g. Ammoniacal Nitrogen, Dissolved
Oxygen, Nitrate) has changed over time, to help spot where water quality is
worst.

Data source: Environment Agency Water Quality Archive
             https://environment.data.gov.uk/water-quality/ (OGL v3)

Run:
    streamlit run app.py

The app reads from a local cache built by `fetch_data.py`. If no cache is
found, it offers a one-click "quick start" fetch of a small sample so you can
see the app working before running the full fetch.
"""

from __future__ import annotations

import io
from datetime import date

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

import fetch_data as fd

# --------------------------------------------------------------------------- #
# 0. PAGE CONFIG
# --------------------------------------------------------------------------- #
st.set_page_config(
    page_title="London Water Quality Explorer",
    page_icon="\U0001F4A7",
    layout="wide",
    initial_sidebar_state="expanded",
)

ONS_COLOURS = {
    "navy": "#003087",
    "aqua_blue": "#27A0CC",
    "dark_blue": "#206095",
    "light_blue": "#A8CEE2",
    "green": "#0F8243",
    "orange": "#F4901E",
    "pink": "#EB4A8A",
    "grey": "#AAAAAA",
    "text": "#222222",
    "grid": "#F0F0F0",
}

# Keyword -> whether a HIGHER reading means WORSE water quality.
# Anything not matched defaults to "higher is worse", the common case for
# pollutant concentrations.
DETERMINAND_DIRECTION = {
    "dissolved oxygen": False,
    "oxygen": False,
    "ph": None,  # neutral-is-best, handled separately
    "ammonia": True,
    "ammoniacal": True,
    "nitrate": True,
    "nitrite": True,
    "phosphate": True,
    "phosphorus": True,
    "bod": True,
    "cod": True,
    "coliform": True,
    "e. coli": True,
    "e.coli": True,
    "enterococci": True,
    "turbidity": True,
    "suspended solids": True,
    "lead": True,
    "copper": True,
    "zinc": True,
    "cadmium": True,
    "mercury": True,
    "arsenic": True,
    "chromium": True,
    "nickel": True,
    "temperature": None,
}


def higher_is_worse(determinand_label: str) -> bool | None:
    label = (determinand_label or "").lower()
    for keyword, direction in DETERMINAND_DIRECTION.items():
        if keyword in label:
            return direction
    return True  # default assumption for unrecognised pollutant-style determinands


# --------------------------------------------------------------------------- #
# 1. DATA LOAD
# --------------------------------------------------------------------------- #
@st.cache_data(show_spinner=False)
def load_cached_data():
    sp, ms = fd.load_cache()
    return sp, ms


@st.cache_data(show_spinner=False)
def quick_start_fetch(radius_km: float, years: int):
    sp = fd.fetch_london_sampling_points(radius_km=radius_km)
    end = date.today()
    start = date(end.year - years, end.month, end.day)
    ms = fd.fetch_measurements(start, end, radius_km=radius_km)
    fd.save_cache(sp, ms)
    return sp, ms


sampling_points, measurements = load_cached_data()

if sampling_points.empty or measurements.empty:
    st.title("London Water Quality Explorer")
    st.warning(
        "No local data cache found yet. Build the full cache from the command line for "
        "best results:\n\n"
        "```\npython fetch_data.py --years 5 --radius-km 30\n```\n\n"
        "Or use the quick-start button below to pull a small sample (15km radius, "
        "2 years) directly from the Environment Agency API so you can try the app now. "
        "This calls the live API and may take a minute or two."
    )
    if st.button("\U0001F680 Quick-start: fetch a small sample now"):
        with st.spinner("Fetching sampling points and measurements from the EA API..."):
            try:
                sampling_points, measurements = quick_start_fetch(radius_km=15, years=2)
            except Exception as exc:  # noqa: BLE001
                st.error(f"Fetch failed: {exc}")
                st.stop()
        st.success(f"Fetched {len(sampling_points)} sampling points and {len(measurements)} measurements.")
        st.rerun()
    st.stop()

# Measurements already carry the sampling point's lat/long/label directly
# (fetched with _view=full), so no merge against the sampling-points table is
# needed for plotting. We still optionally bring in sampling_point_type from
# the sampling-points table, since that's a "list sampling points" property
# rather than a "list measurements" one.
df = measurements.copy()
if "sampling_point_type" not in df.columns and not sampling_points.empty:
    df = df.merge(
        sampling_points[["notation", "sampling_point_type"]],
        left_on="sampling_point_notation",
        right_on="notation",
        how="left",
    )
if "sampling_point_type" not in df.columns:
    df["sampling_point_type"] = pd.NA

df["site_label"] = df["sampling_point_label"].fillna(df["sampling_point_notation"])
df = df.dropna(subset=["lat", "long", "result_numeric", "year"])
df["year"] = df["year"].astype(int)

# --------------------------------------------------------------------------- #
# 2. SIDEBAR FILTERS
# --------------------------------------------------------------------------- #
with st.sidebar:
    st.header("Filters")

    determinand_counts = df["determinand_label"].value_counts()
    determinand_options = determinand_counts.index.tolist()
    default_determinand = next(
        (d for d in determinand_options if "ammonia" in d.lower() or "dissolved oxygen" in d.lower()),
        determinand_options[0] if determinand_options else None,
    )
    determinand = st.selectbox(
        "Determinand (what to measure)",
        options=determinand_options,
        index=determinand_options.index(default_determinand) if default_determinand in determinand_options else 0,
    )

    years_available = sorted(df["year"].unique())
    if len(years_available) > 1:
        year_range = st.select_slider(
            "Year range",
            options=years_available,
            value=(years_available[0], years_available[-1]),
        )
    else:
        year_range = (years_available[0], years_available[0]) if years_available else (None, None)

    sp_types = sorted(df["sampling_point_type"].dropna().unique().tolist())
    selected_types = st.multiselect("Sampling point type", options=sp_types, default=sp_types)

    top_n = st.slider("Number of 'worst' sites to highlight", min_value=3, max_value=20, value=8)

    st.markdown("---")
    st.caption(
        "Data: Environment Agency Water Quality Archive, Open Government Licence v3. "
        "Filters apply to the map, trend chart and ranking table below."
    )

# --------------------------------------------------------------------------- #
# Filtered data
# --------------------------------------------------------------------------- #
mask = (
    (df["determinand_label"] == determinand)
    & (df["year"] >= year_range[0])
    & (df["year"] <= year_range[1])
)
if selected_types:
    mask &= df["sampling_point_type"].isin(selected_types)

fdf = df[mask].copy()

if fdf.empty:
    st.title("London Water Quality Explorer")
    st.info("No measurements match the current filters — try widening the year range or determinand.")
    st.stop()

worse_is_higher = higher_is_worse(determinand)
unit = fdf["unit"].dropna().iloc[0] if fdf["unit"].notna().any() else ""

# Per site-per-year aggregate
site_year = (
    fdf.groupby(["sampling_point_notation", "site_label", "lat", "long", "year"], as_index=False)
    .agg(mean_value=("result_numeric", "mean"), n_samples=("result_numeric", "count"))
)

# A 0-100 "severity score" per year: percentile rank, oriented so 100 = worst
def _severity(group: pd.DataFrame) -> pd.DataFrame:
    pct = group["mean_value"].rank(pct=True) * 100
    group = group.copy()
    if worse_is_higher is False:
        pct = 100 - pct
    group["severity"] = pct
    return group


site_year = site_year.groupby("year", group_keys=False).apply(_severity)

latest_year = int(site_year["year"].max())
latest = site_year[site_year["year"] == latest_year].sort_values("severity", ascending=False)
worst_site_row = latest.iloc[0] if not latest.empty else None

# --------------------------------------------------------------------------- #
# 3. STORY HEADER
# --------------------------------------------------------------------------- #
if worst_site_row is not None:
    st.title(
        f"{worst_site_row['site_label']} shows the worst {determinand.lower()} readings "
        f"in London for {latest_year}"
    )
else:
    st.title(f"London {determinand} water quality, {year_range[0]}–{year_range[1]}")

st.caption(
    f"Comparing {fdf['sampling_point_notation'].nunique()} sampling points across "
    f"{year_range[0]}–{year_range[1]}. "
    + ("Higher values are worse." if worse_is_higher else "Lower values are worse." if worse_is_higher is False else "")
)

# --------------------------------------------------------------------------- #
# 4. KEY METRICS
# --------------------------------------------------------------------------- #
col1, col2, col3, col4 = st.columns(4)
col1.metric("Sampling points", f"{fdf['sampling_point_notation'].nunique():,}")
col2.metric("Measurements", f"{len(fdf):,}")
if worst_site_row is not None:
    col3.metric(
        f"Worst site ({latest_year})",
        f"{worst_site_row['mean_value']:.2f} {unit}",
        help=worst_site_row["site_label"],
    )
best_site_row = latest.iloc[-1] if not latest.empty else None
if best_site_row is not None:
    col4.metric(
        f"Best site ({latest_year})",
        f"{best_site_row['mean_value']:.2f} {unit}",
        help=best_site_row["site_label"],
    )

st.markdown("")

# --------------------------------------------------------------------------- #
# 5. CHARTS
# --------------------------------------------------------------------------- #

def render_chart(fig: go.Figure, chart_id: str, caption: str = ""):
    """Applies ONS/Economist styling and adds a PPTX download button."""
    fig.update_layout(
        font_family="Arial",
        title_font_size=16,
        title_font_color=ONS_COLOURS["text"],
        plot_bgcolor="white",
        paper_bgcolor="white",
        margin=dict(l=40, r=20, t=60, b=40),
        legend=dict(orientation="h", yanchor="bottom", y=1.02),
    )
    st.plotly_chart(fig, use_container_width=True, key=chart_id)
    if caption:
        st.caption(caption)
    _pptx_download_button(fig, chart_id)


def _pptx_download_button(fig: go.Figure, chart_id: str):
    try:
        img_bytes = fig.to_image(format="png", width=1200, height=700, scale=2)
    except Exception:
        st.caption("(Install `kaleido` to enable PPTX export for this chart.)")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.0))

    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.7))
    tf = txbox.text_frame
    tf.text = fig.layout.title.text or chart_id
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.bold = True

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    st.download_button(
        label="\u2b07 Download slide (PPTX)",
        data=buf,
        file_name=f"{chart_id}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key=f"dl_{chart_id}",
    )


tab_map, tab_trend, tab_table = st.tabs(["Map over time", "Trend", "Ranking table"])

with tab_map:
    map_fig = px.scatter_mapbox(
        site_year.sort_values("year"),
        lat="lat",
        lon="long",
        color="severity",
        size="n_samples",
        size_max=18,
        hover_name="site_label",
        hover_data={"mean_value": ":.2f", "n_samples": True, "lat": False, "long": False, "severity": False},
        animation_frame="year",
        color_continuous_scale="Reds",
        range_color=(0, 100),
        mapbox_style="open-street-map",
        zoom=9,
        center={"lat": 51.5074, "lon": -0.1278},
        height=650,
        title=f"{determinand} severity by sampling point, {year_range[0]}–{year_range[1]}",
    )
    map_fig.update_layout(coloraxis_colorbar=dict(title="Severity<br>(100=worst)"))
    render_chart(map_fig, "map_severity_over_time", caption="Drag the year slider below the map, or press play.")

with tab_trend:
    worst_ids = latest.head(top_n)["sampling_point_notation"].tolist()
    trend_df = site_year[site_year["sampling_point_notation"].isin(worst_ids)]

    trend_fig = px.line(
        trend_df.sort_values("year"),
        x="year",
        y="mean_value",
        color="site_label",
        markers=True,
        title=f"The {len(worst_ids)} worst {latest_year} sites for {determinand.lower()}, over time",
        color_discrete_sequence=px.colors.qualitative.Bold,
    )
    trend_fig.update_yaxes(title=f"Mean {determinand} ({unit})", gridcolor=ONS_COLOURS["grid"])
    trend_fig.update_xaxes(title="", showgrid=False)
    render_chart(
        trend_fig,
        "trend_worst_sites",
        caption="Each line is one sampling point, ranked worst-first for the most recent year in range.",
    )

with tab_table:
    show_cols = ["site_label", "year", "mean_value", "n_samples", "severity"]
    ranking = latest[show_cols].rename(
        columns={
            "site_label": "Sampling point",
            "year": "Year",
            "mean_value": f"Mean {determinand} ({unit})",
            "n_samples": "Samples",
            "severity": "Severity (0-100, worst=100)",
        }
    )
    st.dataframe(ranking.reset_index(drop=True), use_container_width=True, height=500)
    st.download_button(
        "\u2b07 Download ranking as CSV",
        ranking.to_csv(index=False).encode("utf-8"),
        file_name=f"london_{determinand.replace(' ', '_').lower()}_ranking_{latest_year}.csv",
        mime="text/csv",
    )

# --------------------------------------------------------------------------- #
# 6. INSIGHT
# --------------------------------------------------------------------------- #
if worst_site_row is not None:
    direction_note = (
        "higher readings indicate worse water quality"
        if worse_is_higher
        else "lower readings indicate worse water quality"
        if worse_is_higher is False
        else "deviation from the healthy range indicates worse water quality"
    )
    st.info(
        f"\U0001F4A1 **Insight:** For {determinand} ({direction_note}), "
        f"**{worst_site_row['site_label']}** was the worst-performing London sampling point in "
        f"{latest_year}, averaging {worst_site_row['mean_value']:.2f} {unit} across "
        f"{int(worst_site_row['n_samples'])} samples. Use the Trend tab to see whether this site "
        f"has been consistently poor or is a recent deterioration."
    )

st.caption(
    "Source: Environment Agency Water Quality Archive, "
    "https://environment.data.gov.uk/water-quality/ — Open Government Licence v3. "
    "© Environment Agency copyright and/or database right."
)
